from PIL import Image
import img2pdf
from dotenv import load_dotenv
from reportlab.lib.pagesizes import letter, A4
from reportlab.pdfgen import canvas
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, Image as RLImage
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT

# Document processing
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import pdfkit

from bot.loger_info.main import logger
from utils.settings import ENV_PATH

load_dotenv(ENV_PATH)


# ============ TO PDF CONVERSION FUNCTIONS ============

async def image_to_pdf(image_path: str, output_path: str):
    """Rasmni PDF ga o'tkazish"""
    try:
        img = Image.open(image_path)
        if img.mode == 'RGBA':
            img = img.convert('RGB')

        pdf_bytes = img2pdf.convert(image_path)
        with open(output_path, "wb") as f:
            f.write(pdf_bytes)
        return True
    except Exception as e:
        logger.error(f"Image to PDF error: {e}")
        return False


async def docx_to_pdf(doc_path: str, output_path: str):
    """Word dokumentni PDF ga o'tkazish"""
    try:
        doc = Document(doc_path)
        pdf = SimpleDocTemplate(output_path, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []

        # Title style
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=16,
            textColor=colors.HexColor('#1a1a1a'),
            spaceAfter=30,
            alignment=TA_CENTER
        )

        # Normal style
        normal_style = ParagraphStyle(
            'CustomNormal',
            parent=styles['Normal'],
            fontSize=11,
            leading=14,
            spaceAfter=12
        )

        for para in doc.paragraphs:
            if para.text.strip():
                # Detect heading vs normal text
                if para.style.name.startswith('Heading'):
                    p = Paragraph(para.text, title_style)
                else:
                    p = Paragraph(para.text, normal_style)
                story.append(p)

        # Add tables if any
        for table in doc.tables:
            data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                data.append(row_data)

            if data:
                t = Table(data)
                t.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 10),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                ]))
                story.append(t)
                story.append(Spacer(1, 20))

        pdf.build(story)
        return True
    except Exception as e:
        logger.error(f"DOCX to PDF error: {e}")
        return False


async def pptx_to_pdf(pptx_path: str, output_path: str):
    """PowerPoint ni PDF ga o'tkazish"""
    try:
        prs = Presentation(pptx_path)
        c = canvas.Canvas(output_path, pagesize=letter)
        width, height = letter

        for slide_num, slide in enumerate(prs.slides, 1):
            # Slide title
            c.setFont("Helvetica-Bold", 18)
            c.drawString(50, height - 50, f"Slide {slide_num}")

            y_position = height - 100
            c.setFont("Helvetica", 12)

            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    # Wrap text
                    words = text.split()
                    line = ""
                    for word in words:
                        if len(line + word) < 80:
                            line += word + " "
                        else:
                            c.drawString(50, y_position, line)
                            y_position -= 20
                            line = word + " "
                            if y_position < 50:
                                c.showPage()
                                y_position = height - 50

                    if line:
                        c.drawString(50, y_position, line)
                        y_position -= 30

                    if y_position < 50:
                        c.showPage()
                        y_position = height - 50

            # New page for next slide
            if slide_num < len(prs.slides):
                c.showPage()

        c.save()
        return True
    except Exception as e:
        logger.error(f"PPTX to PDF error: {e}")
        return False


async def xlsx_to_pdf(excel_path: str, output_path: str):
    """Excel faylni PDF ga o'tkazish"""
    try:
        wb = load_workbook(excel_path)
        pdf = SimpleDocTemplate(output_path, pagesize=A4)
        elements = []
        styles = getSampleStyleSheet()

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]

            # Sheet title
            title = Paragraph(f"<b>{sheet_name}</b>", styles['Heading1'])
            elements.append(title)
            elements.append(Spacer(1, 20))

            # Get data
            data = []
            for row in ws.iter_rows(values_only=True):
                row_data = [str(cell) if cell is not None else "" for cell in row]
                data.append(row_data)

            if data:
                # Create table
                col_widths = [1.2 * inch] * len(data[0]) if data[0] else []
                table = Table(data, colWidths=col_widths)
                table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 9),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 10),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black),
                    ('FONTSIZE', (0, 1), (-1, -1), 8),
                ]))
                elements.append(table)

            elements.append(PageBreak())

        pdf.build(elements)
        return True
    except Exception as e:
        logger.error(f"XLSX to PDF error: {e}")
        return False


async def html_to_pdf(html_path: str, output_path: str):
    """HTML ni PDF ga o'tkazish"""
    try:
        # pdfkit ishlatish (wkhtmltopdf kerak)
        # Alternativa: weasyprint ishlatish mumkin
        options = {
            'page-size': 'A4',
            'encoding': 'UTF-8',
            'enable-local-file-access': None
        }
        pdfkit.from_file(html_path, output_path, options=options)
        return True
    except Exception as e:
        logger.error(f"HTML to PDF error: {e}")
        # Fallback: oddiy usul
        try:
            with open(html_path, 'r', encoding='utf-8') as f:
                html_content = f.read()

            from reportlab.platypus import Paragraph
            from reportlab.lib.styles import getSampleStyleSheet

            pdf = SimpleDocTemplate(output_path, pagesize=letter)
            styles = getSampleStyleSheet()
            story = []

            # HTML teglarini olib tashlash
            import re
            clean_text = re.sub('<[^<]+?>', '', html_content)

            for line in clean_text.split('\n'):
                if line.strip():
                    p = Paragraph(line.strip(), styles['Normal'])
                    story.append(p)
                    story.append(Spacer(1, 12))

            pdf.build(story)
            return True
        except Exception as e2:
            logger.error(f"HTML to PDF fallback error: {e2}")
            return False
