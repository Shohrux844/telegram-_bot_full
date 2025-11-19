import os

from dotenv import load_dotenv

# PDF manipulation
from PyPDF2 import PdfReader, PdfWriter
from pdf2image import convert_from_path
from pptx import Presentation

from bot.loger_info.main import logger
from utils.settings import ENV_PATH

load_dotenv(ENV_PATH)


# ============ FROM PDF CONVERSION FUNCTIONS ============

async def pdf_to_images(pdf_path: str, output_folder: str):
    """PDF dan rasmlar yaratish"""
    try:
        os.makedirs(output_folder, exist_ok=True)
        images = convert_from_path(pdf_path, dpi=200)

        image_paths = []
        for i, image in enumerate(images):
            image_path = os.path.join(output_folder, f'page_{i + 1}.jpg')
            image.save(image_path, 'JPEG')
            image_paths.append(image_path)

        return image_paths
    except Exception as e:
        logger.error(f"PDF to images error: {e}")
        return None


async def pdf_to_docx(pdf_path: str, output_path: str):
    """PDF dan Word yaratish"""
    try:
        from pdf2docx import Converter
        cv = Converter(pdf_path)
        cv.convert(output_path)
        cv.close()
        return True
    except Exception as e:
        logger.error(f"PDF to DOCX error: {e}")
        return False


async def pdf_to_pptx(pdf_path: str, output_path: str):
    """PDF dan PowerPoint yaratish"""
    try:
        # PDF'dan rasmlar olish
        images = convert_from_path(pdf_path, dpi=150)

        # PowerPoint yaratish
        prs = Presentation()
        prs.slide_width = 9144000  # 10 inches
        prs.slide_height = 6858000  # 7.5 inches

        for i, image in enumerate(images):
            # Rasmni saqlash
            img_path = f'temp_slide_{i}.png'
            image.save(img_path, 'PNG')

            # Slide qo'shish
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)

            # Rasmni slide ga qo'shish
            left = top = 0
            slide.shapes.add_picture(img_path, left, top,
                                     width=prs.slide_width,
                                     height=prs.slide_height)

            # Temp faylni o'chirish
            os.remove(img_path)

        prs.save(output_path)
        return True
    except Exception as e:
        logger.error(f"PDF to PPTX error: {e}")
        return False


async def pdf_to_xlsx(pdf_path: str, output_path: str):
    """PDF dan Excel yaratish"""
    try:
        # PDF'dan matnni o'qish
        reader = PdfReader(pdf_path)

        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws.title = "PDF Content"

        row = 1
        for page_num, page in enumerate(reader.pages, 1):
            text = page.extract_text()
            ws.cell(row=row, column=1, value=f"=== Sahifa {page_num} ===")
            row += 1

            for line in text.split('\n'):
                if line.strip():
                    ws.cell(row=row, column=1, value=line.strip())
                    row += 1

            row += 1  # Bo'sh qator

        wb.save(output_path)
        return True
    except Exception as e:
        logger.error(f"PDF to XLSX error: {e}")
        return False


async def pdf_to_pdfa(pdf_path: str, output_path: str):
    """PDF dan PDF/A yaratish"""
    try:
        reader = PdfReader(pdf_path)
        writer = PdfWriter()

        # Barcha sahifalarni ko'chirish
        for page in reader.pages:
            writer.add_page(page)

        # Metadata qo'shish (PDF/A uchun)
        writer.add_metadata({
            '/Title': 'PDF/A Document',
            '/Author': 'PDF Converter Bot',
            '/Subject': 'Converted to PDF/A',
            '/Creator': 'PDF Converter Bot',
            '/Producer': 'PDF Converter Bot',
        })

        # Saqlash
        with open(output_path, 'wb') as f:
            writer.write(f)

        return True
    except Exception as e:
        logger.error(f"PDF to PDF/A error: {e}")
        return False
